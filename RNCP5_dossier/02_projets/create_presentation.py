#!/usr/bin/env python3
"""Génère la présentation PPTX pour le passage oral RNCP 5 DWWM — FloraShop"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
import os

SCREENSHOTS = '/root/holbertonschool-web_back_end/RNCP5_dossier/02_projets/screenshots'
OUTPUT = '/root/holbertonschool-web_back_end/RNCP5_dossier/02_projets/Presentation_RNCP5_Jaille_Dimitri.pptx'

PINK = RGBColor(0xBC, 0x62, 0x88)
PINK_DARK = RGBColor(0xAB, 0x59, 0x7C)
PINK_LIGHT = RGBColor(0xF5, 0xE0, 0xE8)
BG_CREAM = RGBColor(0xFF, 0xF6, 0xF2)
TEXT_DARK = RGBColor(0x23, 0x11, 0x11)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height
BAR_H = Inches(0.45)
BODY_TOP = Inches(1.2)
BODY_MAX_H = Inches(5.5)


def add_bg(slide, color=BG_CREAM):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_bar(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), SH - BAR_H, SW, BAR_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = PINK
    bar.line.fill.background()
    tf = bar.text_frame
    tf.text = "Jaille Dimitri — FloraShop | Pivoine & Lilas — RNCP 5 DWWM — Holberton School — 2026"
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def ttl(slide, text, top=Inches(0.25), left=Inches(0.5), width=None, size=Pt(28), color=PINK_DARK):
    if width is None:
        width = SW - Inches(1)
    tx = slide.shapes.add_textbox(left, top, width, Inches(0.6))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = True
    p.font.color.rgb = color


def bdy(slide, text, top=BODY_TOP, left=Inches(0.6), width=None,
        height=BODY_MAX_H, size=Pt(14), color=TEXT_DARK,
        bold=False, align=PP_ALIGN.LEFT):
    if width is None:
        width = SW - Inches(1.2)
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = size
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
        p.space_after = Pt(2)
        p.space_before = Pt(0)


def pic(slide, name, left, top, width=None, height=None):
    path = os.path.join(SCREENSHOTS, name)
    if os.path.exists(path):
        kw = {}
        if width: kw['width'] = width
        if height: kw['height'] = height
        slide.shapes.add_picture(path, left, top, **kw)


# === 1 — PAGE DE GARDE ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, PINK)
pic(s, 'holberton_logo.png', Inches(5.9), Inches(0.3), width=Inches(1.5))
ttl(s, "FloraShop — Pivoine & Lilas", top=Inches(1.8), size=Pt(40), color=WHITE)
bdy(s, "Application e-commerce pour artisan fleuriste",
    top=Inches(2.6), size=Pt(22), color=PINK_LIGHT)
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(3.2), Inches(4), Inches(0.04))
ln.fill.solid(); ln.fill.fore_color.rgb = WHITE; ln.line.fill.background()
bdy(s,
    "Jaille Dimitri\n"
    "Co-développeur : Mattieu Mourroux\n\n"
    "Titre visé : Développeur Web et Web Mobile (DWWM)\n"
    "RNCP Niveau 5 — Code RNCP 37674\n"
    "Holberton School — Session 2025-2026 — Mars 2026",
    top=Inches(3.5), size=Pt(16), color=WHITE, height=Inches(3))
pic(s, 'accueil_desktop.png', Inches(8.5), Inches(1.5), width=Inches(4.3))


# === 2 — SOMMAIRE ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); ttl(s, "Sommaire"); add_bar(s)
bdy(s, '\n'.join([
    "1.  Contexte du projet et cahier des charges",
    "2.  Stack technique et environnement",
    "3.  Architecture de l'application (3-tiers)",
    "4.  Front-end : interfaces statiques et dynamiques",
    "5.  Back-end : base de données (MCD, MLD, SQL, ORM)",
    "6.  Back-end : API REST et composants métier",
    "7.  Sécurité de l'application",
    "8.  Démonstration et jeu d'essai",
    "9.  Difficultés rencontrées et enseignements",
    "10. Conclusion et perspectives",
]), size=Pt(18))


# === 3 — CONTEXTE ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); ttl(s, "1. Contexte du projet"); add_bar(s)
bdy(s,
    "Commanditaire : Pivoine & Lilas\n"
    "→ Artisan fleuriste, gérée par la mère de Mattieu Mourroux\n\n"
    "Objectif : concevoir une boutique en ligne complète\n"
    "→ Vente de fleurs, compositions, vases, parfums, cadeaux\n"
    "→ Paiement sécurisé en ligne (Stripe)\n"
    "→ Administration : produits, catégories, utilisateurs\n\n"
    "Type : Projet en entreprise\n"
    "Formation : DWWM — Holberton School\n\n"
    "Équipe :\n"
    "• Jaille Dimitri — Développeur\n"
    "• Mattieu Mourroux — Co-développeur",
    size=Pt(16), width=Inches(7))


# === 4 — CAHIER DES CHARGES ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); ttl(s, "Cahier des charges — Besoins fonctionnels"); add_bar(s)
bdy(s, '\n'.join([
    "• Catalogue produits par catégories (7) avec filtres et recherche",
    "• Panier d'achat dynamique (ajout, suppression, calcul total)",
    "• Paiement en ligne sécurisé via Stripe (EUR)",
    "• Gestion des comptes utilisateurs (inscription, connexion, JWT)",
    "• Tableau d'administration CRUD (produits, catégories, utilisateurs)",
    "• API REST documentée via Swagger UI (Flask-RESTX)",
    "• Pages vitrine : accueil, portfolio, événementiel, entreprises, contact",
    "• Abonnement floral avec paiement récurrent",
    "",
    "Besoins non-fonctionnels :",
    "• Interface responsive (desktop, tablette, mobile)",
    "• Sécurité : hashage, JWT, validation des entrées",
    "• Architecture MVC, pattern Repository, code documenté",
]), size=Pt(15))


# === 5 — STACK ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); ttl(s, "2. Stack technique"); add_bar(s)
bdy(s,
    "Front-end\n"
    "• HTML5, CSS3 (charte vintage #bc6288), JS vanilla\n"
    "• Jinja2, Stripe.js v3, Google reCAPTCHA\n"
    "• 16 pages HTML, 4 fichiers CSS (932 lignes)\n\n"
    "Back-end\n"
    "• Python 3.10, Flask 3.1.0, Flask-RESTX 1.3.0\n"
    "• Flask-SQLAlchemy 3.1.1, Flask-Migrate 4.1.0\n"
    "• PyJWT, Werkzeug, Flask-CORS",
    top=BODY_TOP, size=Pt(15), width=Inches(6))
bdy(s,
    "Base de données\n"
    "• PostgreSQL — 7 tables, ORM SQLAlchemy\n\n"
    "Paiement\n"
    "• Stripe API (Payment Intents + Webhooks)\n\n"
    "Outils\n"
    "• Git / GitHub, python-dotenv, coverage",
    top=BODY_TOP, left=Inches(7), size=Pt(15), width=Inches(5.5))


# === 6 — ARCHITECTURE ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); ttl(s, "3. Architecture 3-tiers"); add_bar(s)
bdy(s,
    "┌──────────────────────────────────────┐\n"
    "│  NAVIGATEUR (Présentation)           │\n"
    "│  HTML5 / CSS3 / JS / Stripe.js       │\n"
    "└───────────────┬──────────────────────┘\n"
    "                │ HTTP / JSON\n"
    "┌───────────────▼──────────────────────┐\n"
    "│  SERVEUR FLASK (Logique métier)      │\n"
    "│  Routes → Services → Repositories    │\n"
    "│  Flask-RESTX │ Facade │ StripeService │\n"
    "│  Auth JWT │ Swagger │ Factory         │\n"
    "└───────────────┬──────────────────────┘\n"
    "                │ SQL (psycopg2)\n"
    "┌───────────────▼──────────────────────┐\n"
    "│  POSTGRESQL (Données)                │\n"
    "│  7 tables — ORM SQLAlchemy           │\n"
    "└──────────────────────────────────────┘",
    top=Inches(1.1), size=Pt(13), left=Inches(1), width=Inches(6.5))
bdy(s,
    "Design Patterns :\n\n"
    "• MVC (Model-View-Controller)\n"
    "• Repository (accès aux données)\n"
    "• Facade (point d'accès unique)\n"
    "• Factory (create_app())",
    top=Inches(1.3), left=Inches(8.5), size=Pt(15), width=Inches(4), color=PINK_DARK)


# === 7 — COMPÉTENCES ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); ttl(s, "Compétences DWWM couvertes"); add_bar(s)
bdy(s,
    "AT1 — Front-end sécurisé\n\n"
    "✅ CP1 — Maquetter des interfaces\n"
    "    12 maquettes, design responsive\n\n"
    "✅ CP2 — Interfaces statiques\n"
    "    16 templates HTML5/CSS3, charte vintage\n\n"
    "✅ CP3 — Interfaces dynamiques\n"
    "    JS vanilla (api.js, auth.js), panier, Stripe.js",
    size=Pt(15), width=Inches(6))
bdy(s,
    "AT2 — Back-end sécurisé\n\n"
    "✅ CP4 — Base de données relationnelle\n"
    "    PostgreSQL, 7 tables, 6 migrations\n\n"
    "✅ CP5 — Composants d'accès aux données\n"
    "    Pattern Repository, ORM SQLAlchemy, Facade\n\n"
    "✅ CP6 — Composants métier serveur\n"
    "    REST API Swagger, StripeService, JWT, rôles",
    left=Inches(6.8), size=Pt(15), width=Inches(6))


# === 8 — SCREENSHOTS DESKTOP ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); ttl(s, "4. Front-end — Pages principales (desktop)"); add_bar(s)
for i, (f, lbl) in enumerate([
    ('accueil_desktop.png', 'Accueil'), ('boutique_desktop.png', 'Boutique'),
    ('panier_desktop.png', 'Panier'), ('checkout_desktop.png', 'Checkout')
]):
    x = Inches(0.3 + i * 3.2)
    bdy(s, lbl, top=Inches(1.0), left=x, width=Inches(3), size=Pt(13),
        color=PINK_DARK, bold=True, align=PP_ALIGN.CENTER, height=Inches(0.3))
    pic(s, f, x, Inches(1.4), width=Inches(3))


# === 9 — SCREENSHOTS MOBILE ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); ttl(s, "Front-end — Version mobile (responsive)"); add_bar(s)
for i, (f, lbl) in enumerate([
    ('accueil_mobile.png', 'Accueil'), ('boutique_mobile.png', 'Boutique'),
    ('panier_mobile.png', 'Panier'), ('checkout_mobile.png', 'Checkout')
]):
    x = Inches(0.8 + i * 3.1)
    bdy(s, lbl, top=Inches(1.0), left=x, width=Inches(2.5), size=Pt(13),
        color=PINK_DARK, bold=True, align=PP_ALIGN.CENTER, height=Inches(0.3))
    pic(s, f, x, Inches(1.4), height=Inches(5))


# === 10 — CSS / RESPONSIVE ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); ttl(s, "Front-end — Charte graphique et responsive"); add_bar(s)
bdy(s,
    "Charte graphique — Variables CSS\n\n"
    ":root {\n"
    "  --vintage-bg: #fff6f2;\n"
    "  --vintage-main: #bc6288;\n"
    "  --vintage-text: #231111;\n"
    "}\n\n"
    "Polices : EB Garamond + Inter\n"
    "Navigation fixe, bordure rose 2px",
    size=Pt(14), width=Inches(5.5))
bdy(s,
    "Responsive Design\n\n"
    "• <meta viewport> sur toutes les pages\n"
    "• Media query 900px : sidebar empilée\n"
    "• Media query 600px : grille 1 colonne\n"
    "• Navigation compacte sur mobile\n"
    "• Grille : auto-fill minmax(260px, 1fr)\n\n"
    "Accessibilité\n"
    "• Attributs alt sur les images\n"
    "• Contrastes suffisants\n"
    "• Labels, required, type=email",
    left=Inches(6.5), size=Pt(14), width=Inches(6))


# === 11 — JS DYNAMIQUE ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); ttl(s, "Front-end — Partie dynamique (JavaScript)"); add_bar(s)
bdy(s,
    "ApiService (api.js — 261 lignes)\n"
    "• Appels API REST (fetch + async/await)\n"
    "• Token JWT dans localStorage\n"
    "• Header Authorization: Bearer auto\n"
    "• login(), register(), logout(), getProducts()…\n\n"
    "Panier (panier.html)\n"
    "• Stockage localStorage (JSON)\n"
    "• renderCart() : reconstruction DOM\n"
    "• updateCartCount() : badge temps réel\n\n"
    "Stripe.js v3 (checkout.html)\n"
    "• Stripe Elements : iframe sécurisé (PCI DSS)\n"
    "• client_secret depuis le serveur\n"
    "• Données bancaires jamais sur notre serveur",
    size=Pt(15))


# === 12 — MCD ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); ttl(s, "5. Base de données — MCD Merise"); add_bar(s)
bdy(s,
    "7 entités, 6 associations\n\n"
    "USERS ──(PASSER 0,N / 0,1)──→ ORDERS\n"
    "USERS ──(RÉDIGER 0,N / 1,1)──→ REVIEWS\n"
    "ORDERS ──(COMPOSER 1,N / 1,1)──→ ORDER_ITEMS\n"
    "ORDER_ITEMS ──(RÉFÉRENCER 1,1 / 0,N)──→ PRODUCTS\n"
    "CATEGORIES ──(CONTENIR 1,N / 1,1)──→ PRODUCTS\n"
    "PRODUCTS ──(HISTORISER 0,N / 1,1)──→ PRICES",
    size=Pt(14), width=Inches(7))
bdy(s,
    "Règles de gestion clés\n\n"
    "• RG1 : (username, email) unique\n"
    "• RG3 : Prix figé au moment de l'achat\n"
    "• RG4 : Suppression catégorie → cascade\n"
    "• RG6 : Historique prix (is_active)\n"
    "• RG7 : pending → paid | failed | cancelled",
    left=Inches(7.5), size=Pt(14), width=Inches(5))


# === 13 — SQL + ORM ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); ttl(s, "Base de données — SQL et ORM SQLAlchemy"); add_bar(s)
bdy(s,
    "7 tables PostgreSQL (SERIAL PK)\n\n"
    "users        → id, username, email, password, is_admin\n"
    "categories   → id, name, created_at\n"
    "products     → id, name, price, category_id (FK), is_on_sale\n"
    "orders       → id, user_id (FK NULL), email, total, stripe_id, status\n"
    "order_items  → id, order_id (FK CASCADE), product_id, qty, price\n"
    "reviews      → id, content, rating, user_id (FK)\n"
    "prices       → id, amount, is_active, product_id (FK)",
    size=Pt(14), width=Inches(7))
bdy(s,
    "Correspondance SQL ↔ ORM\n\n"
    "• SERIAL → primary_key=True\n"
    "• REFERENCES → db.ForeignKey()\n"
    "• CASCADE → cascade='all, delete-orphan'\n"
    "• to_dict() exclut le password\n\n"
    "6 migrations Alembic",
    left=Inches(7.5), size=Pt(14), width=Inches(5))


# === 14 — API REST ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); ttl(s, "6. API REST — Flask-RESTX et Swagger"); add_bar(s)
bdy(s,
    "Endpoints par namespace :\n\n"
    "Auth       POST /auth/login, /auth/register\n"
    "Products   GET / POST / PUT / DELETE\n"
    "Categories GET / POST / PUT / DELETE\n"
    "Users      GET / POST / PUT / DELETE\n"
    "Reviews    GET / POST / PATCH / DELETE\n"
    "Payments   POST /create-payment-intent, /confirm, /webhook\n\n"
    "• Doc auto-générée : localhost:5000/api/v1\n"
    "• Auth Bearer JWT dans Swagger\n"
    "• Factory pattern : create_app()",
    size=Pt(14), width=Inches(6.5))
pic(s, 'swagger_desktop.png', Inches(7.5), Inches(1.2), width=Inches(5.3))


# === 15 — COMPOSANTS MÉTIER ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); ttl(s, "Composants métier côté serveur"); add_bar(s)
bdy(s,
    "Pattern Repository\n"
    "• CategoryRepo, ProductRepo, UserRepo, ReviewRepo, PriceRepo\n"
    "• CRUD : get_all, get_by_id, create, update, delete\n"
    "• Validation unicité, vérification ID, rollback\n"
    "• ORM → protection injection SQL\n\n"
    "Pattern Facade\n"
    "• Point d'accès unique à tous les repositories\n"
    "• Simplifie l'utilisation depuis les routes\n\n"
    "StripeService\n"
    "• Total calculé depuis les prix en BDD\n"
    "• Création Payment Intent + commande\n"
    "• Vérification webhook (signature crypto)\n"
    "• Gestion transactionnelle : flush → commit / rollback",
    size=Pt(15))


# === 16 — SÉCURITÉ ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); ttl(s, "7. Sécurité de l'application"); add_bar(s)
bdy(s,
    "Mesures implémentées ✅\n\n"
    "• Hashage mots de passe (Werkzeug)\n"
    "• Authentification JWT (PyJWT, exp. 1 jour)\n"
    "• CORS restrictif (localhost)\n"
    "• reCAPTCHA sur inscription\n"
    "• Validation entrées (champs, types, unicité)\n"
    "• ORM SQLAlchemy (anti-injection SQL)\n"
    "• Variables sensibles dans .env\n"
    "• Webhook Stripe signé\n"
    "• Cascade delete\n"
    "• Échappement Jinja2 (anti-XSS)\n"
    "• to_dict() exclut le password",
    size=Pt(14), width=Inches(6))
bdy(s,
    "Veille sécurité — OWASP Top 10\n\n"
    "Sources : OWASP, Flask Security,\n"
    "Stripe Security, CVE Details, Snyk\n\n"
    "12 failles identifiées :\n"
    "F1 : MdP en clair (auth.py)\n"
    "F2 : MdP admin hardcodé\n"
    "F3 : Clé JWT hardcodée\n"
    "F5 : DEBUG=True + 0.0.0.0\n"
    "F6 : MdP exposé via GET\n"
    "F7 : Endpoints sans auth\n"
    "F8 : CORS trop permissif\n"
    "F12 : Logs de MdP en clair\n\n"
    "→ Corrections proposées dans le dossier",
    left=Inches(7), size=Pt(14), width=Inches(5.5))


# === 17 — JEU D'ESSAI ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); ttl(s, "8. Jeu d'essai — Processus de commande"); add_bar(s)
bdy(s,
    "7 scénarios testés :\n\n"
    "✅ Test 1 : GET /products → 200, 6 produits\n"
    "✅ Test 2 : GET /categories → 200, 4 catégories\n"
    "✅ Test 3 : POST /auth/login → 200, token JWT\n"
    "⚠️  Test 4 : POST /create-payment-intent → timeout (clés test)\n"
    "⚠️  Test 5 : POST /confirm-payment → dépend du test 4\n"
    "✅ Test 6 : POST /categories → 201, ID auto-généré\n"
    "✅ Test 7 : DELETE /categories/6 → 200, cascade OK\n\n"
    "Résultat : 5/7 tests validés\n"
    "2 tests attendus en échec (clés Stripe non configurées)",
    size=Pt(15))


# === 18 — DIFFICULTÉS ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); ttl(s, "9. Difficultés et enseignements"); add_bar(s)
bdy(s,
    "Difficultés techniques\n\n"
    "• Intégration Stripe : Payment Intents, centimes, webhooks\n"
    "  → Solution : StripeService dédié\n\n"
    "• Architecture Repository : séparation métier / données\n"
    "  → Solution : Facade + flush/commit/rollback\n\n"
    "• Migrations Alembic : 6 versions, schéma en évolution\n"
    "  → Rigueur dans le versioning",
    size=Pt(15), width=Inches(6))
bdy(s,
    "Enseignements tirés\n\n"
    "• Importance des design patterns\n"
    "  (MVC, Repository, Facade)\n\n"
    "• La sécurité dès la conception\n\n"
    "• Intégration API tierces :\n"
    "  idempotence, webhooks, erreurs async\n\n"
    "• Documentation automatisée\n"
    "  (Swagger) toujours à jour",
    left=Inches(7), size=Pt(15), width=Inches(5.5))


# === 19 — CONCLUSION ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); ttl(s, "10. Conclusion et perspectives"); add_bar(s)
bdy(s,
    "Bilan\n\n"
    "• Application e-commerce full-stack fonctionnelle\n"
    "• 6 compétences DWWM couvertes (CP1 à CP6)\n"
    "• Architecture : MVC, Repository, Facade, Factory\n"
    "• Audit de sécurité avec analyse OWASP",
    size=Pt(15), width=Inches(6))
bdy(s,
    "Perspectives d'évolution\n\n"
    "• Déploiement : Gunicorn + Nginx, SSL\n"
    "• Sécurité : hashage systématique, rotation JWT,\n"
    "  rate limiting\n"
    "• Fonctionnalités : emails, stocks, avis,\n"
    "  analytics commerçant\n"
    "• Tests : pytest + coverage\n"
    "• Performance : Redis, pagination, eager loading",
    left=Inches(7), size=Pt(15), width=Inches(5.5))


# === 20 — MERCI ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, PINK)
ttl(s, "Merci de votre attention", top=Inches(2.3), size=Pt(40), color=WHITE)
bdy(s, "Questions ?", top=Inches(3.2), size=Pt(28), color=PINK_LIGHT)
ln2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(4.0), Inches(3), Inches(0.04))
ln2.fill.solid(); ln2.fill.fore_color.rgb = WHITE; ln2.line.fill.background()
bdy(s,
    "Jaille Dimitri — Mattieu Mourroux\n"
    "FloraShop — Pivoine & Lilas\n"
    "Holberton School — Session 2025-2026",
    top=Inches(4.3), size=Pt(16), color=WHITE, height=Inches(2))
pic(s, 'holberton_logo.png', Inches(10), Inches(5), width=Inches(2))


prs.save(OUTPUT)
print(f"Présentation créée : {OUTPUT}")
print(f"Nombre de slides : {len(prs.slides)}")
